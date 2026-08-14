from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import os

def test_native_pptx_chart():
    print("⏳ 正在生成原生可编辑 PPTX 商业研报...")
    
    # 1. 初始化 PPT 文档
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # 添加一张带标题的空白幻灯片
    
    # 2. 设置标题
    title = slide.shapes.title
    title.text = "Eco-Research 宏观商业分析报告 (Demo)"

    # 3. 模拟 Agent 清洗出来的结构化 JSON 数据
    chart_data = CategoryChartData()
    chart_data.categories = ['比亚迪 (BYD)', '特斯拉 (Tesla)', '长城 (GWM)', '吉利 (Geely)']
    chart_data.add_series('2026年马来西亚销量预测 (万辆)', (4.2, 2.1, 1.5, 1.8))

    # 4. 在 PPT 中插入原生的柱状图 (完美支持双击唤出 Excel)
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    
    # 5. 保存到项目根目录
    current_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(current_dir)
    output_filename = os.path.join(project_root, 'Eco_Research_Demo.pptx')
    
    prs.save(output_filename)
    print(f"🎉 成功！请在左侧目录查看生成的：Eco_Research_Demo.pptx")

if __name__ == "__main__":
    test_native_pptx_chart()