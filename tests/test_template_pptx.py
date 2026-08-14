import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def get_available_templates(template_dir):
    """读取文件夹，获取所有可用的 PPT 模板"""
    if not os.path.exists(template_dir):
        print(f"⚠️ 找不到模板文件夹: {template_dir}")
        return []
    
    # 筛选出所有的 .pptx 文件
    templates = [f for f in os.listdir(template_dir) if f.endswith('.pptx')]
    return templates

def generate_report_with_template(template_name):
    """使用指定的模板生成报告"""
    # 1. 拼凑模板的绝对路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(current_dir)
    template_path = os.path.join(project_root, 'templates', 'ppt', template_name)
    
    print(f"⏳ 正在加载模板: {template_name} ...")
    
    # 2. 核心魔法：直接把模板路径传进去！系统会继承它的所有样式！
    prs = Presentation(template_path)
    
    # 3. 添加一页新幻灯片 (使用模板自带的空白版式，通常索引是 5 或 6)
    # 注意：不同模板的母版索引可能不同，如果报错，可以尝试把 5 换成 1 或 6
    try:
        slide_layout = prs.slide_layouts[5] 
    except IndexError:
        slide_layout = prs.slide_layouts[0] # 如果没有6个版式，就用默认第1个
        
    slide = prs.slides.add_slide(slide_layout)
    
    # 4. 插入图表数据 (大模型分析的结果)
    chart_data = CategoryChartData()
    chart_data.categories = ['比亚迪', '特斯拉', '长城', '吉利']
    chart_data.add_series('2026年销量预测 (万辆)', (4.2, 2.1, 1.5, 1.8))

    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    
    # 5. 保存生成的报告
    output_filename = os.path.join(project_root, f'最终报告_{template_name}')
    prs.save(output_filename)
    print(f"🎉 报告已成功套用模板并生成: {output_filename}")

if __name__ == "__main__":
    # 定义模板文件夹路径 (假设你已经把 ppt模板 改名为了 templates/ppt)
    current_dir = os.path.dirname(os.path.abspath(__file__))
    root_dir = os.path.dirname(current_dir)
    template_folder = os.path.join(root_dir, 'templates', 'ppt')
    
    # 获取模板列表
    available_templates = get_available_templates(template_folder)
    
    if available_templates:
        print("👇 发现以下精美模板：")
        for i, t in enumerate(available_templates):
            print(f"[{i+1}] {t}")
        
        # 模拟前端用户的下拉框选择（这里我们默认选第1个测试）
        # 实际你在做网页端时，这里就是前端传过来的值
        chosen_template = available_templates[0] 
        print(f"\n👉 模拟用户选择了: {chosen_template}")
        
        # 执行生成
        generate_report_with_template(chosen_template)