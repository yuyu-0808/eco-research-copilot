import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def test_auto_render_ppt():
    print("⏳ 正在启动 Eco-Research PPT 渲染引擎...")
    
    # 1. 这里就是原本需要大模型生成的 JSON（分镜脚本）
    # 我们先手动模拟一份完美的数据，测试渲染引擎
    mock_llm_json = {
        "slides": [
            {
                "slide_type": "title_page",
                "title": "马来西亚新能源出海报告",
                "subtitle": "2026 宏观经济与市场格局分析 | Eco-Research Copilot 自动生成"
            },
            {
                "slide_type": "content_page",
                "title": "核心宏观结论 (PESTEL)",
                "bullets": [
                    "政策端 (P)：马来西亚政府对纯电动车 (EV) 免征进口税及消费税，政策红利期显著。",
                    "经济端 (E)：人均可支配收入稳步增长，中产阶级对 15-20 万区间的新能源车接受度极高。",
                    "基建端 (T)：目前全国公共充电桩仅约 2000 个，充电焦虑是最大的市场痛点。"
                ]
            },
            {
                "slide_type": "chart_page",
                "title": "2026 本地市场销量预测",
                "chart_data": {"比亚迪": 4.2, "特斯拉": 2.1, "长城": 1.5, "吉利": 1.8}
            }
        ]
    }

    # 2. 定位你的“纯净版”模板
    current_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(current_dir)
    template_path = os.path.join(project_root, 'templates', 'ppt', 'clean_template.pptx')
    output_path = os.path.join(project_root, '最终完美版_马来西亚行研.pptx')

    if not os.path.exists(template_path):
        print(f"❌ 找不到模板文件，请确保路径正确: {template_path}")
        return

    # 3. 开始向纯净模板里“填空”
    prs = Presentation(template_path)

    for slide_data in mock_llm_json["slides"]:
        # 渲染封面页
        if slide_data["slide_type"] == "title_page":
            slide_layout = prs.slide_layouts[0] # 母版第 0 个通常是封面
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = slide_data["subtitle"]
                
        # 渲染内容页
        elif slide_data["slide_type"] == "content_page":
            slide_layout = prs.slide_layouts[1] # 母版第 1 个通常是内容页
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]
            
            # 填入项目符号
            tf = slide.placeholders[1].text_frame
            tf.text = slide_data["bullets"][0]
            for bullet in slide_data["bullets"][1:]:
                p = tf.add_paragraph()
                p.text = bullet
                
        # 渲染图表页
        elif slide_data["slide_type"] == "chart_page":
            # 找一个尽量空白的版式（索引可能是 5 或 6，这里如果报错，可以把 6 改成 5 或 1）
            try:
                slide_layout = prs.slide_layouts[6] 
            except IndexError:
                slide_layout = prs.slide_layouts[1] 
                
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["title"]
            
            # 画出原生图表
            chart_data = CategoryChartData()
            chart_data.categories = list(slide_data["chart_data"].keys())
            chart_data.add_series('预测销量 (万辆)', tuple(slide_data["chart_data"].values()))
            
            x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
            slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    # 4. 保存输出
    prs.save(output_path)
    print(f"🎉 成功！报告已生成，完全没有宇宙飞船重叠了！请查看: {output_path}")

if __name__ == "__main__":
    test_auto_render_ppt()